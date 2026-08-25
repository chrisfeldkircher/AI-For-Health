"""Build the ML4Health mid-term talk on the TUM HLU template.

Run from project root with the datascience env:
  & "$env:USERPROFILE\AppData\Local\miniconda3\envs\datascience\python.exe" presentation/build_pptx.py

Structure: what we used, what went wrong, why, what we change next, references.
Plain wording (no dashes as punctuation, no slogans). Figures from presentation/figures/.

NOTE on author names: the title slide has Christoph Feldkircher plus two visible
placeholders ("___ Lee", "___ Chouksey") because the other first names were not
provided. Replace the underscores in PowerPoint, or tell me the names.
"""
from __future__ import annotations
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "HLU_Presentation_Template_Oct_2022_Modified_by_Tian.pptx"
FIG = ROOT / "presentation" / "figures"
OUT = ROOT / "presentation" / "MidTerm_Postmortem_FINAL_PLAN.pptx"

CLASS_LINE = "ML4Health SS26, Technical University of Munich"
AUTHOR_LINE = "Christoph Feldkircher, ___ Lee, ___ Chouksey"
FOOTER = "Feldkircher, Lee, Chouksey  |  ML4Health SS26, TUM  |  URTIC Cold Detection"

NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _strip_slides(prs):
    lst = prs.slides._sldIdLst
    rids = [s.get(f"{{{NS_R}}}id") for s in list(lst)]
    for s in list(lst):
        lst.remove(s)
    for rid in rids:
        try:
            prs.part.drop_rel(rid)
        except KeyError:
            pass


def _layout(prs, master_idx, name):
    for lay in prs.slide_masters[master_idx].slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(name)


def _ph_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _ph_name(slide, sub):
    for ph in slide.placeholders:
        if sub.lower() in ph.name.lower():
            return ph
    return None


def _set_text(tf, text, size=None, bold=False, color=None):
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    if size is not None:
        r.font.size = Pt(size)
    if bold:
        r.font.bold = True
    if color is not None:
        r.font.color.rgb = RGBColor(*color)


def _footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.34), Inches(5.31), Inches(8.6), Inches(0.30))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    _set_text(tf, text, size=9, color=(110, 110, 110))


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _picture(slide, image_path):
    pic_ph = None
    for ph in slide.placeholders:
        if "PICTURE" in str(ph.placeholder_format.type):
            pic_ph = ph
            break
    if pic_ph is not None:
        sp = pic_ph._element
        sp.getparent().remove(sp)
    with Image.open(str(image_path)) as im:
        w, h = im.size
    aspect = w / h
    top, max_h, max_w = 1.05, 4.05, 9.5
    hh = max_w / aspect
    if hh <= max_h:
        ww, hh = max_w, hh
    else:
        hh, ww = max_h, max_h * aspect
    left = (10.0 - ww) / 2.0
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top),
                             width=Inches(ww), height=Inches(hh))


def title_slide(prs, lay, title, lines, notes):
    s = prs.slides.add_slide(lay)
    t = _ph_idx(s, 0)
    if t is not None:
        _set_text(t.text_frame, title, size=24, bold=True, color=(0, 51, 89))
    c = _ph_idx(s, 10)
    if c is not None:
        c.text_frame.clear()
        for i, ln in enumerate(lines):
            p = c.text_frame.paragraphs[0] if i == 0 else c.text_frame.add_paragraph()
            r = p.add_run(); r.text = ln
            r.font.size = Pt(15 if i == 0 else 13)
            r.font.color.rgb = RGBColor(60, 60, 60)
    _footer(s, FOOTER)
    _notes(s, notes)


def image_slide(prs, lay, title, image_path, notes):
    s = prs.slides.add_slide(lay)
    t = _ph_idx(s, 0)
    if t is not None:
        _set_text(t.text_frame, title, size=22, bold=True, color=(0, 0, 0))
    _picture(s, image_path)
    _footer(s, FOOTER)
    _notes(s, notes)


def bullets_slide(prs, lay, title, bullets, notes, size=15):
    s = prs.slides.add_slide(lay)
    t = _ph_idx(s, 0)
    if t is not None:
        _set_text(t.text_frame, title, size=22, bold=True, color=(0, 0, 0))
    c = _ph_name(s, "inhaltsplatzhalter") or _ph_idx(s, 1)
    if c is not None:
        c.text_frame.clear()
        for i, b in enumerate(bullets):
            p = c.text_frame.paragraphs[0] if i == 0 else c.text_frame.add_paragraph()
            p.level = 0
            r = p.add_run(); r.text = b
            r.font.size = Pt(size)
            r.font.color.rgb = RGBColor(20, 20, 20)
    _footer(s, FOOTER)
    _notes(s, notes)


def main():
    prs = Presentation(str(TEMPLATE))
    _strip_slides(prs)
    L_title = _layout(prs, 0, "1_Start")
    L_img = _layout(prs, 4, "1_große Bilder")
    L_bul = _layout(prs, 4, "4_Inhalt + Text")

    title_slide(
        prs, L_title,
        "URTIC Cold Detection: mid-term result and plan",
        [AUTHOR_LINE, "", CLASS_LINE, "June 2026"],
        "We submitted our system to the ComParE 2017 Cold task. It scored UAR 0.62 on the "
        "hidden test, below what our own validation predicted. I will cover what we used, "
        "what went wrong, why, and what we change next.",
    )

    image_slide(
        prs, L_img, "What we used", FIG / "fig6_architecture.png",
        "We rebuilt the 2017 late fusion with a modern backbone. Frozen WavLM with "
        "audit-weighted pooling, two handcrafted groups fused on top, five seeds averaged, "
        "one threshold tuned on a dev split. On that split we matched the published baseline.",
    )

    image_slide(
        prs, L_img, "What went wrong", FIG / "fig1_confusion_matrix.png",
        "On the hidden test we got 0.62. Our own validation said 0.69, and the result sits "
        "below every dev fold we measured. Precision on cold is 14 percent. We checked the run "
        "five ways and it is byte for byte correct, so this is not a bug. It is a generalization gap.",
    )

    image_slide(
        prs, L_img, "The problem is the ranking, not the threshold",
        FIG / "fig8_recall_below_shadow.png",
        "The easy story is that we call cold too often. But at our threshold the cold-call rate "
        "is about 42 percent on dev too, so the rate is matched. Hold it matched and look at "
        "recall: on test it is below every dev fold, on both classes. At a fixed call rate that "
        "is a ranking loss, and no threshold change fixes a ranking loss.",
    )

    image_slide(
        prs, L_img, "Why our dev folds did not catch it",
        FIG / "fig10_pool_shift.png",
        "Our folds are speaker-disjoint, but they all come from one speaker pool. The hidden "
        "test is a different population. The cold direction we learned leans on speaker traits "
        "from our pool and transfers worse. Folds drawn from the same pool cannot measure that gap.",
    )

    image_slide(
        prs, L_img, "What we change next", FIG / "fig9_what_we_change.png",
        "We ruled out the cheap fixes by measurement. Two real options. One, report this as the "
        "ceiling of a frozen backbone, which our protocol already flagged. Two, unfreeze the top "
        "layers and train with speaker-stratified batches so the cold cue depends less on who is "
        "talking. We pick on the dev folds first.",
    )

    bullets_slide(
        prs, L_bul, "References",
        [
            "Schuller et al. The INTERSPEECH 2017 Computational Paralinguistics Challenge: "
            "Addressee, Cold and Snoring. Interspeech 2017.",
            "Chen et al. WavLM: Large-Scale Self-Supervised Pre-Training for Full Stack Speech "
            "Processing. IEEE JSTSP, 2022.",
            "Yang et al. SUPERB: Speech processing Universal PERformance Benchmark. Interspeech 2021.",
            "Coppock, Jones, Kiskin, Schuller. COVID-19 detection from audio: seven grains of salt. "
            "Lancet Digital Health, 2021.",
        ],
        "References for the baseline, the backbone, the frozen-representation paradigm, and the "
        "speaker-confound point.",
        size=13,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"[done] wrote {OUT} ({OUT.stat().st_size/1024:.1f} KB, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
